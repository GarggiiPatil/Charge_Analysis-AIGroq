from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import os

# ---------- SAFE ABSOLUTE PATH ----------
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_PATH = os.path.join(BASE_DIR, "assets", "TATA_LOGO.jpg")

def generate_report_ppt(report: dict) -> BytesIO:
    prs = Presentation()

    # ---------- HEADER + FOOTER ----------
    def add_header_footer(slide, page_no):
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(
                LOGO_PATH,
                Inches(0.2), Inches(0.15),
                height=Inches(0.6)
            )

        footer = slide.shapes.add_textbox(
            Inches(8.5), Inches(6.8),
            Inches(1), Inches(0.3)
        )
        p = footer.text_frame.paragraphs[0]
        p.text = str(page_no)
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.RIGHT

    # ---------- TITLE SLIDE ----------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Charge Analysis Report"
    slide.placeholders[1].text = "Generated using AI"
    add_header_footer(slide, 1)

    # ---------- SUMMARY SLIDE ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for key, value in report["summary"].items():
        p = tf.add_paragraph()
        p.text = f"{key.replace('_', ' ').title()}: {value}"
        p.font.size = Pt(14)

    add_header_footer(slide, 2)

    # ---------- INSIGHTS SLIDE ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Insights"
    slide.placeholders[1].text = report["insights"]
    add_header_footer(slide, 3)

    # ---------- TABLE SLIDE ----------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "City-wise Report"

    table_data = report["report_table"]
    rows = len(table_data) + 1
    cols = len(table_data[0])

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(4)
    ).table

    # Table header
    for i, col in enumerate(table_data[0].keys()):
        table.cell(0, i).text = col

    # Table rows
    for r, row in enumerate(table_data, start=1):
        for c, val in enumerate(row.values()):
            table.cell(r, c).text = str(val)

    add_header_footer(slide, 4)

    # ---------- EXPORT ----------
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer